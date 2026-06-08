"""
Insert nof_fig_flow.png into slide 5 of MiniTelcoPlatform_v2.pptx.
Replaces the '[Architecture Diagram]' placeholder with the actual image.
"""
from pptx import Presentation
from pptx.util import Emu
from PIL import Image
import os, sys

PPTX_IN  = '/mnt/c/Users/turka/Desktop/MiniTelcoPlatform_v2.pptx'
PPTX_OUT = '/mnt/c/Users/turka/Desktop/MiniTelcoPlatform_v2.pptx'
IMG_PATH = '/home/turkad/mini-telco-platform/docs/nof_fig_flow.png'

prs = Presentation(PPTX_IN)
slide = prs.slides[4]          # 0-indexed → slide 5

# ── 1. Find & remove the placeholder shape ───────────────────────────────────
ph_left = ph_top = ph_width = ph_height = None
to_remove = []
for shape in slide.shapes:
    if hasattr(shape, 'text') and '[ Architecture Diagram ]' in shape.text:
        ph_left   = shape.left
        ph_top    = shape.top
        ph_width  = shape.width
        ph_height = shape.height
        to_remove.append(shape._element)
        print(f"Placeholder found: left={ph_left}, top={ph_top}, "
              f"cx={ph_width}, cy={ph_height}")

if not to_remove:
    print("ERROR: placeholder not found — is the text exactly '[ Architecture Diagram ]'?")
    sys.exit(1)

for el in to_remove:
    el.getparent().remove(el)

# ── 2. Compute image dimensions (maintain aspect ratio, fill height) ──────────
with Image.open(IMG_PATH) as im:
    px_w, px_h = im.size
print(f"Image pixels: {px_w} × {px_h}  ratio={px_w/px_h:.4f}")

# Fit inside the placeholder box — constrain by height first, then check width
img_h = ph_height
img_w = int(img_h * px_w / px_h)

if img_w > ph_width:
    # too wide — constrain by width instead
    img_w = ph_width
    img_h = int(img_w * px_h / px_w)

# Centre within the placeholder area
img_left = ph_left + (ph_width  - img_w) // 2
img_top  = ph_top  + (ph_height - img_h) // 2

print(f"Placing image: left={img_left} top={img_top} cx={img_w} cy={img_h}")

# ── 3. Insert the image ───────────────────────────────────────────────────────
slide.shapes.add_picture(IMG_PATH,
                         Emu(img_left), Emu(img_top),
                         Emu(img_w),    Emu(img_h))

prs.save(PPTX_OUT)
print(f"OK: saved → {PPTX_OUT}")
